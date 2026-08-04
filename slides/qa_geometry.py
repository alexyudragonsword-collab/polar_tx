"""Geometry QA for the deck: bounds, margins, text overflow, overlaps.

Substitute for pixel QA (LibreOffice cannot render in this sandbox).
Text-fit is estimated from character widths: CJK glyphs are ~1 em wide,
Latin ~0.55 em, which is what makes a Chinese deck overflow where an
English one fits.
"""
import math
import sys
import unicodedata

from pptx import Presentation
from pptx.util import Emu

SLIDE_W, SLIDE_H = 13.333, 7.5
MARGIN_MIN = 0.5


def _em(v):
    return Emu(v).inches if v is not None else None


def _wide(ch):
    return unicodedata.east_asian_width(ch) in ("W", "F")


def _text_width_em(t):
    return sum(1.0 if _wide(c) else 0.55 for c in t)


def audit(path):
    prs = Presentation(path)
    problems = []
    for idx, slide in enumerate(prs.slides, 1):
        boxes = []
        for sh in slide.shapes:
            try:
                x, y = _em(sh.left), _em(sh.top)
                w, h = _em(sh.width), _em(sh.height)
            except Exception:
                continue
            if None in (x, y, w, h):
                continue
            name = sh.shape_type
            # 1) out of bounds
            if x < -0.05 or y < -0.05 or x + w > SLIDE_W + 0.05 \
                    or y + h > SLIDE_H + 0.05:
                # decorative circles are deliberately bled off-slide
                if not (sh.shape_type is not None and w > 2 and h > 2
                        and not sh.has_text_frame):
                    problems.append(
                        f"S{idx}: OUT OF BOUNDS {name} "
                        f"x={x:.2f} y={y:.2f} w={w:.2f} h={h:.2f}")
            # 2) text overflow estimate
            if sh.has_text_frame and sh.text_frame.text.strip():
                for para in sh.text_frame.paragraphs:
                    txt = "".join(r.text for r in para.runs)
                    if not txt.strip():
                        continue
                    sz = None
                    for r in para.runs:
                        if r.font.size:
                            sz = r.font.size.pt
                            break
                    if sz is None:
                        continue
                    usable_w = max(w - 0.1, 0.2)
                    cap = usable_w * 72.0 / sz          # em per line
                    lines = max(1, math.ceil(_text_width_em(txt) / cap))
                    need = lines * sz * 1.22 / 72.0
                    if need > h + 0.02:
                        problems.append(
                            f"S{idx}: TEXT OVERFLOW ({need:.2f}\" needed vs "
                            f"{h:.2f}\" box, {lines} lines @ {sz:.0f}pt): "
                            f"{txt[:42]!r}")
            if sh.has_text_frame and sh.text_frame.text.strip():
                boxes.append((x, y, w, h, sh.text_frame.text[:26]))
        # 3) overlapping TEXT boxes (cards behind text are fine)
        for i in range(len(boxes)):
            for j in range(i + 1, len(boxes)):
                a, b = boxes[i], boxes[j]
                ox = min(a[0] + a[2], b[0] + b[2]) - max(a[0], b[0])
                oy = min(a[1] + a[3], b[1] + b[3]) - max(a[1], b[1])
                if ox > 0.06 and oy > 0.06:
                    problems.append(
                        f"S{idx}: TEXT OVERLAP {ox:.2f}x{oy:.2f}\" "
                        f"{a[4]!r} <-> {b[4]!r}")
    return problems, len(prs.slides)


if __name__ == "__main__":
    probs, n = audit(sys.argv[1])
    print(f"{n} slides audited")
    if not probs:
        print("PASS: no geometry problems")
    else:
        print(f"{len(probs)} issue(s):")
        for p in probs:
            print("  " + p)
